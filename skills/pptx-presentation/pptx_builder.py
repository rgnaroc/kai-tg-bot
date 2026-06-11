#!/usr/bin/env python3
"""
PPTX Presentation Builder — skill для Kai AI Assistant
Генерирует .pptx из JSON-описания.
"""

import json
import sys
import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────── Defaults ───────────────────────────────

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

DEFAULT_THEME = {
    "primary_color": "1B3A5C",
    "accent_color": "2196F3",
    "bg_color": "FFFFFF",
    "text_color": "333333",
    "font_name": "Roboto",
    "font_size_title": 36,
    "font_size_body": 18,
    "font_size_code": 14,
}

# ─────────────────────────────── Helpers ───────────────────────────────

def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def merge_theme(base, overrides):
    theme = dict(base)
    if overrides:
        theme.update(overrides)
    return theme

def add_text_box(slide, left, top, width, height, text, font_name="Roboto",
                 font_size=18, bold=False, color="333333", alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color)
    p.alignment = alignment
    return txBox

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)

def add_accent_bar(slide, theme):
    """Горизонтальная полоска accent вверху слайда"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Pt(6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme["accent_color"])
    shape.line.fill.background()

def add_page_number(slide, num, total, theme):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_name=theme["font_name"],
                 font_size=10, color="999999", alignment=PP_ALIGN.RIGHT)

# ─────────────────────────────── Slide Builders ───────────────────────────────

def build_title_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, theme["primary_color"])

    # accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(3.2), Inches(2.0), Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme["accent_color"])
    shape.line.fill.background()

    # title
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(1.5),
                 data.get("title", "Presentation"),
                 font_name=theme["font_name"], font_size=44, bold=True,
                 color="FFFFFF", alignment=PP_ALIGN.LEFT)

    # subtitle
    if data.get("subtitle"):
        add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.8),
                     data["subtitle"],
                     font_name=theme["font_name"], font_size=22, bold=False,
                     color="B0C4DE", alignment=PP_ALIGN.LEFT)

    # author + date at bottom
    bottom_text = data.get("author", "")
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
                 bottom_text,
                 font_name=theme["font_name"], font_size=12,
                 color="8FA8C8", alignment=PP_ALIGN.LEFT)

    add_page_number(slide, page_num, total, theme)


def build_bullets_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["bg_color"])
    add_accent_bar(slide, theme)

    # title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_title"],
                 bold=True, color=theme["primary_color"])

    # separator line under title
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.2), Inches(2.0), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(theme["accent_color"])
    shape.line.fill.background()

    # bullet items
    items = data.get("items", [])
    body_text = "\n" + "\n".join(f"• {item}" for item in items)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0),
                 body_text,
                 font_name=theme["font_name"], font_size=theme["font_size_body"],
                 color=theme["text_color"])

    add_page_number(slide, page_num, total, theme)


def build_table_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["bg_color"])
    add_accent_bar(slide, theme)

    # title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_title"],
                 bold=True, color=theme["primary_color"])

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if not headers or not rows:
        return

    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    tbl_left = Inches(0.8)
    tbl_top = Inches(1.6)
    tbl_width = Inches(11.5)
    tbl_height = Inches(0.6) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table

    # header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = hex_to_rgb("FFFFFF")
            paragraph.font.name = theme["font_name"]
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(theme["primary_color"])
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = hex_to_rgb(theme["text_color"])
                paragraph.font.name = theme["font_name"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb("F5F7FA") if r_idx % 2 == 0 else hex_to_rgb("FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_page_number(slide, page_num, total, theme)


def build_image_text_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["bg_color"])
    add_accent_bar(slide, theme)

    # title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_title"],
                 bold=True, color=theme["primary_color"])

    # text on the left
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0),
                 data.get("text", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_body"],
                 color=theme["text_color"])

    # image on the right (if path provided)
    img_path = data.get("image_path", "")
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(7.0), Inches(1.6),
                                     width=Inches(5.5))
        except Exception:
            add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(1.0),
                         "[Image could not be loaded]",
                         font_name=theme["font_name"], font_size=12, color="999999")
    else:
        # placeholder rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb("E8EDF2")
        shape.line.fill.background()
        add_text_box(slide, Inches(7.0), Inches(3.8), Inches(5.5), Inches(0.6),
                     "[Изображение]",
                     font_name=theme["font_name"], font_size=16, color="999999",
                     alignment=PP_ALIGN.CENTER)

    add_page_number(slide, page_num, total, theme)


def build_process_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["bg_color"])
    add_accent_bar(slide, theme)

    # title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_title"],
                 bold=True, color=theme["primary_color"])

    steps = data.get("steps", [])
    step_text = ""
    for i, step in enumerate(steps, 1):
        step_text += f"{i}. {step}\n\n"

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0),
                 step_text,
                 font_name=theme["font_name"], font_size=theme["font_size_body"],
                 color=theme["text_color"])

    add_page_number(slide, page_num, total, theme)


def build_code_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["bg_color"])
    add_accent_bar(slide, theme)

    # title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=theme["font_size_title"],
                 bold=True, color=theme["primary_color"])

    # code block background
    code = data.get("code", "")
    lang = data.get("language", "")

    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = hex_to_rgb("1E1E2E")
    code_bg.line.fill.background()

    # language label
    if lang:
        add_text_box(slide, Inches(1.2), Inches(1.8), Inches(3), Inches(0.4),
                     f"▸ {lang}",
                     font_name=theme["font_name"], font_size=11, bold=True,
                     color="888888")

    # code content
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10.5), Inches(4.4),
                 code,
                 font_name="Consolas", font_size=theme["font_size_code"],
                 color="F8F8F2")

    add_page_number(slide, page_num, total, theme)


def build_quote_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "F5F7FA")
    add_accent_bar(slide, theme)

    quote_text = data.get("text", "")
    source = data.get("source", "")

    # quote mark
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(1), Inches(1.0),
                 "❝",
                 font_name=theme["font_name"], font_size=60, bold=True,
                 color=theme["accent_color"], alignment=PP_ALIGN.LEFT)

    # quote body
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.5), Inches(3.5),
                 quote_text,
                 font_name=theme["font_name"], font_size=28, bold=False,
                 color=theme["primary_color"])

    # source
    if source:
        add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
                     f"— {source}",
                     font_name=theme["font_name"], font_size=16, bold=False,
                     color="888888")

    add_page_number(slide, page_num, total, theme)


def build_section_slide(prs, data, theme, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, theme["primary_color"])

    # big centered title
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(2.0),
                 data.get("title", ""),
                 font_name=theme["font_name"], font_size=40, bold=True,
                 color="FFFFFF", alignment=PP_ALIGN.CENTER)

    add_page_number(slide, page_num, total, theme)


# ─────────────────────────────── Main Builder ───────────────────────────────

SLIDE_BUILDERS = {
    "title": build_title_slide,
    "bullets": build_bullets_slide,
    "table": build_table_slide,
    "image_text": build_image_text_slide,
    "process": build_process_slide,
    "code": build_code_slide,
    "quote": build_quote_slide,
    "section": build_section_slide,
}

def build_presentation(data, output_path="output.pptx"):
    theme = merge_theme(DEFAULT_THEME, data.get("theme", {}))

    # Convert string hex colors
    for key in ["primary_color", "accent_color", "bg_color", "text_color"]:
        if isinstance(theme.get(key), str):
            pass  # keep as string, hex_to_rgb called per use

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = data.get("slides", [])
    total = len(slides_data)

    # If no slides, auto-add title slide
    if total == 0:
        slides_data = [{"type": "title"}]

    for i, slide_data in enumerate(slides_data, 1):
        slide_type = slide_data.get("type", "bullets")
        builder = SLIDE_BUILDERS.get(slide_type, build_bullets_slide)
        builder(prs, slide_data, theme, i, total)

    prs.save(output_path)
    return output_path


# ─────────────────────────────── CLI ───────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="PPTX Presentation Builder")
    parser.add_argument("--output", "-o", default="presentation.pptx",
                        help="Output file path")
    parser.add_argument("--input", "-i", help="JSON input file (default: stdin)")
    args = parser.parse_args()

    if args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        raw = sys.stdin.read()
        if not raw.strip():
            print("❌ No input data. Pass JSON via stdin or --input file", file=sys.stderr)
            sys.exit(1)
        data = json.loads(raw)

    path = build_presentation(data, args.output)
    print(f"✅ PPTX created: {path}", file=sys.stderr)
    print(path)


if __name__ == "__main__":
    main()
